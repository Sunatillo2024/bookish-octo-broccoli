import os
from pathlib import Path
import uuid
from pptx import Presentation
from pptx.util import Inches, Pt
from app.models import SlideType, SlideContent, PresentationRequest
from app.config import settings

class PPTGenerator:
    def __init__(self):
        # Ensure storage directory exists
        os.makedirs(settings.STORAGE_PATH, exist_ok=True)

    def generate_presentation(self, request: PresentationRequest) -> str:
        """Generate a PowerPoint slide based on the request"""
        prs = Presentation()

        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = request.title
        subtitle.text = f"By {request.author}"

        # Add content slides
        for slide_content in request.slides:
            self._add_slide(prs, slide_content)

        # Save the presentation
        file_id = str(uuid.uuid4())
        file_path = os.path.join(settings.STORAGE_PATH, f"{file_id}.pptx")
        prs.save(file_path)

        return file_path

    def _add_slide(self, prs: Presentation, content: SlideContent):
        """Add a slide based on its type and content"""
        if content.type == SlideType.TITLE:
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            title.text = content.title
            if content.content:
                subtitle.text = content.content

        elif content.type == SlideType.CONTENT:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = content.title
            if content.content:
                body.text = content.content

        elif content.type == SlideType.BULLET_POINTS:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            body = slide.placeholders[1]
            title.text = content.title

            if content.bullet_points:
                tf = body.text_frame
                tf.text = ""  # Clear default text

                for point in content.bullet_points:
                    p = tf.add_paragraph()
                    p.text = point
                    p.level = 0

        elif content.type == SlideType.TWO_COLUMN:
            slide_layout = prs.slide_layouts[3]  # Assuming layout 3 is two-content
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = content.title

            # Handle columns - this may vary based on your pptx template
            left = slide.placeholders[1]
            right = slide.placeholders[2]

            if content.column1:
                left.text = content.column1
            if content.column2:
                right.text = content.column2

        elif content.type == SlideType.IMAGE:
            # Basic image slide
            slide_layout = prs.slide_layouts[5]  # Blank slide with title
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            title.text = content.title

            # Note: In a real application, you would handle image downloads
            # and insertion here. For simplicity, we're omitting this.